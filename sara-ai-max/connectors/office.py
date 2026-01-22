"""
Office Automation - Word, Excel, PowerPoint control.

This module provides Microsoft Office automation capabilities.
"""

import logging
from typing import Optional, List
from pathlib import Path

logger = logging.getLogger(__name__)


class OfficeAutomation:
    """Automate Microsoft Office applications."""
    
    def __init__(self):
        """Initialize office automation."""
        self.word_available = False
        self.excel_available = False
        self.powerpoint_available = False
        
        # Check for required libraries
        try:
            import docx
            self.word_available = True
        except ImportError:
            logger.warning("python-docx not available. Word automation unavailable.")
        
        try:
            import openpyxl
            self.excel_available = True
        except ImportError:
            logger.warning("openpyxl not available. Excel automation unavailable.")
        
        try:
            from pptx import Presentation
            self.powerpoint_available = True
        except ImportError:
            logger.warning("python-pptx not available. PowerPoint unavailable.")
    
    def create_word_document(
        self,
        filename: str,
        content: List[str],
        save_path: Optional[str] = None
    ) -> bool:
        """
        Create a Word document.
        
        Args:
            filename: Document filename
            content: List of paragraphs
            save_path: Optional save location
            
        Returns:
            True if successful
        """
        if not self.word_available:
            logger.error("Word automation not available")
            return False
        
        try:
            from docx import Document
            
            doc = Document()
            for paragraph in content:
                doc.add_paragraph(paragraph)
            
            if not save_path:
                save_path = str(Path.home() / "Documents" / filename)
            
            doc.save(save_path)
            logger.info(f"Word document created: {save_path}")
            return True
            
        except Exception as e:
            logger.error(f"Error creating Word document: {e}")
            return False
    
    def create_excel_spreadsheet(
        self,
        filename: str,
        data: List[List],
        save_path: Optional[str] = None
    ) -> bool:
        """
        Create an Excel spreadsheet.
        
        Args:
            filename: Spreadsheet filename
            data: 2D list of cell data
            save_path: Optional save location
            
        Returns:
            True if successful
        """
        if not self.excel_available:
            logger.error("Excel automation not available")
            return False
        
        try:
            from openpyxl import Workbook
            
            wb = Workbook()
            ws = wb.active
            
            for row in data:
                ws.append(row)
            
            if not save_path:
                save_path = str(Path.home() / "Documents" / filename)
            
            wb.save(save_path)
            logger.info(f"Excel spreadsheet created: {save_path}")
            return True
            
        except Exception as e:
            logger.error(f"Error creating Excel spreadsheet: {e}")
            return False
    
    def create_powerpoint(
        self,
        filename: str,
        slides: List[dict],
        save_path: Optional[str] = None
    ) -> bool:
        """
        Create a PowerPoint presentation.
        
        Args:
            filename: Presentation filename
            slides: List of slide dictionaries with 'title' and 'content'
            save_path: Optional save location
            
        Returns:
            True if successful
        """
        if not self.powerpoint_available:
            logger.error("PowerPoint automation not available")
            return False
        
        try:
            from pptx import Presentation
            from pptx.util import Inches
            
            prs = Presentation()
            
            for slide_data in slides:
                slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and content layout
                slide.shapes.title.text = slide_data.get('title', '')
                
                if 'content' in slide_data:
                    content = slide.placeholders[1]
                    content.text = slide_data['content']
            
            if not save_path:
                save_path = str(Path.home() / "Documents" / filename)
            
            prs.save(save_path)
            logger.info(f"PowerPoint presentation created: {save_path}")
            return True
            
        except Exception as e:
            logger.error(f"Error creating PowerPoint: {e}")
            return False
